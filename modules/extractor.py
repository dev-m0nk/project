import os
from pptx import Presentation
from PIL import Image
import pytesseract

def extract_slide_texts_and_images(pptx_path, use_ocr=False):
    prs = Presentation(pptx_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        images_data = []
        notes_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            if shape.shape_type == 13:  # Picture
                image = shape.image
                img_bytes = image.blob
                img_path = f"slide_{idx}_image.png"
                with open(img_path, "wb") as f:
                    f.write(img_bytes)

                ocr_text = ""
                if use_ocr:
                    try:
                        img = Image.open(img_path)
                        ocr_text = pytesseract.image_to_string(img)
                    except Exception as e:
                        ocr_text = f"[OCR error: {e}]"

                images_data.append({"path": img_path, "ocr_text": ocr_text})

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            "index": idx,
            "texts": texts,
            "images": images_data,
            "notes": notes_text
        })

    return slides_data
